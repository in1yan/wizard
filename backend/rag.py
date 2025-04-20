import os
import uuid
from typing import List
from dotenv import load_dotenv
from langchain_core.documents import Document
from langchain_core.messages import SystemMessage
from langchain_core.tools import tool
from langchain_groq import ChatGroq
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEmbeddings
from pathlib import Path
from langchain_community.document_loaders import PyPDFLoader, DirectoryLoader
from langgraph.graph import START, StateGraph, MessagesState, END
from langgraph.prebuilt import ToolNode, tools_condition
from langgraph.checkpoint.memory import MemorySaver
from youtube_transcript_api import YouTubeTranscriptApi
from yt_dlp import YoutubeDL
from pptx import Presentation
import docx
load_dotenv()

class ConversationalRAG:
    def __init__(self):
        self.groq_api_key = os.getenv('GROQ_API_KEY')
        self.thread = uuid.uuid4()
        if not self.groq_api_key:
            raise ValueError('GROQ_API_KEY environment variable not set')

        # Initialize the LLM
        self.llm = ChatGroq(
            api_key=self.groq_api_key,
            model_name="meta-llama/llama-4-scout-17b-16e-instruct",
            temperature=0.3
        )

        # Initialize embeddings
        self.embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-mpnet-base-v2")

        # Initialize text splitter
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )
        self.memory = MemorySaver()
        # Initialize vector store
        self.vector_store = None

        # Initialize the graph
        self.graph = None

    def load_documents(self, path: str) -> List[Document]:
        """Load documents from a file or directory"""
        path = Path(path)
        if not path.exists():
            raise FileNotFoundError(f"Path not found: {path}")
        
        documents = []
        if path.is_file() and path.suffix.lower() == '.pdf':
            # Load single PDF file
            loader = PyPDFLoader(str(path))
            documents.extend(loader.load())
        elif path.is_dir():
            # Load all PDFs from directory
            loader = DirectoryLoader(
                str(path),
                glob="**/*.pdf",
                loader_cls=PyPDFLoader
            )
            documents.extend(loader.load())
        elif path.is_file() and path.suffix.lower() in (".ppt", ".pptx"):
            loader = self.ppt_loader(str(path))
            documents.extend( loader )
        elif path.is_file() and path.suffix.lower() in (".doc", ".docx") :
            loader = self.doc_loader(str(path))
            documents.extend( loader )
        else:
            raise ValueError(f"Unsupported file type: {path}")
        
        return documents

    def process_documents(self, documents: List[Document]):
        """Process documents and create vector store"""
        # Split documents into chunks
        splits = self.text_splitter.split_documents(documents)

        # Create vector store
        if self.vector_store is None:
            self.vector_store = FAISS.from_documents(
                documents=splits,
                embedding=self.embeddings
            )
        else:
            # Add new documents to existing vector store
            self.vector_store.add_documents(splits)

        # Set up the graph after documents are processed
        self._setup_graph()
    def ppt_loader(self, path:str):
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        text.append(run.text)
        documents = [Document(page_content="".join(text))]
        return documents
    def doc_loader(self, path:str):
        dox = docx.Document(path)
        paras = dox.paragraphs
        text = ""
        for para in paras:
            text += para.text
        documents = [Document(page_content=text)]
        return documents
    def load_youtube(self, id: str) -> List[Document]:
        """Load documents from a youtube video"""
        transcript = YouTubeTranscriptApi.get_transcript(id)
        text = ""
        for i in transcript:
            text += i['text'] + " "
        documents = [Document(page_content=text)]
        with YoutubeDL() as ydl:
            info_dict = ydl.extract_info(id, download=False)
            title = info_dict.get('title', None)
        return documents, title
    def process_youtube(self, documents: List[Document]) -> str:
        """Process documents and create vector store"""
        # Split documents into chunks
        splits = self.text_splitter.split_documents(documents)

        # Create vector stor
        if self.vector_store is None:
            self.vector_store = FAISS.from_documents(
                documents=splits,
                embedding=self.embeddings
            )
        else:
            # Add new documents to existing vector store
            self.vector_store.add_documents(splits)

        # Set up the graph after documents are processed
        self._setup_graph()
    def _setup_graph(self):
        """Set up the LangGraph components"""
        # Define the retrieval tool
        @tool(response_format="content_and_artifact")
        def retrieve(query: str):
            """Retrieve information related to a query."""
            docs = self.vector_store.similarity_search(query, k=2)
            serialized = "\n\n".join(
                (f"Source: {doc.metadata}\n" f"Content: {doc.page_content}") for doc in docs
            )
            return serialized, docs

        # Node to use the tool if necessary
        def query_or_respond(state: MessagesState):
            """Generate tool call for retrieval or respond."""
            llm_tools = self.llm.bind_tools([retrieve])
            response = llm_tools.invoke(state["messages"])
            return {"messages": [response]}

        # Create tool node
        tools = ToolNode([retrieve])

        # Node to generate the final answer
        def generate(state: MessagesState) -> MessagesState:
            recent_tool_msgs = []
            for message in reversed(state["messages"]):
                if message.type == "tool":
                    recent_tool_msgs.append(message)
                else:
                    break
            tool_msgs = recent_tool_msgs[::-1]
            docs_content = "\n\n".join(doc.content for doc in tool_msgs)

            system_message_content = (
                "You are an assistant for question-answering tasks. "
                "Use the following pieces of retrieved context to answer "
                "the question. If you don't know the answer, say that you don't know. "
                "Use three sentences maximum and keep the answer concise.\n\n"
                f"{docs_content}"
            )

            conversation = [
                msg for msg in state["messages"]
                if msg.type in ("human", "system") or (msg.type == "ai" and not msg.tool_calls)
            ]
            prompt = [SystemMessage(system_message_content)] + conversation
            response = self.llm.invoke(prompt)

            return {"messages": [response]}

        # Construct the graph
        graph_builder = StateGraph(MessagesState)
        graph_builder.add_node( query_or_respond)
        graph_builder.add_node( tools)
        graph_builder.add_node( generate)

        graph_builder.set_entry_point("query_or_respond")
        graph_builder.add_conditional_edges(
            "query_or_respond",
            tools_condition,
            {END: END, "tools": "tools"}
        )
        graph_builder.add_edge("tools", "generate")
        graph_builder.add_edge("generate", END)

        self.graph = graph_builder.compile(checkpointer=self.memory)

    def chat(self, question: str) -> str:
        """Process a question and return the response"""
        if not self.vector_store:
            raise ValueError("No documents have been processed. Please process documents first.")
        if not self.graph:
            raise ValueError("Graph not initialized. Please process documents first.")
        config = {"configurable": {"thread_id": f"{self.thread}"}}
        state = self.graph.invoke({"messages": [{"role": "user", "content": question}]}, config=config)
        return state["messages"][-1].content

# Example usage
if __name__ == "__main__":
    # Initialize the RAG system
    rag = ConversationalRAG()

    # Load and process documents (can be a single PDF or a directory)
    documents = rag.load_documents("./uploads")
    rag.process_documents(documents)

    while True:
        question = input("\nEnter your question (or 'quit' to exit): ")
        if question.lower() == 'quit':
            break
        
        response = rag.chat(question)
        print("\nResponse:", response)
