from langchain_core.documents import Document
from youtube_transcript_api import YouTubeTranscriptApi
from yt_dlp import YoutubeDL
from pptx import Presentation
import docx
from typing import List


def ppt_loader(path:str):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    text.append(run.text)
    documents = [Document(page_content="".join(text), meta_data={"file-type":"ppt", "file-name":path.split('/')[-1]})]
    return documents


def doc_loader( path:str):
    dox = docx.Document(path)
    paras = dox.paragraphs
    text = ""
    for para in paras:
        text += para.text
    documents = [Document(page_content=text, meta_data={"file-type":"docx", "file-name":path.split('/')[-1]})]
    return documents


def load_youtube( id: str) -> List[Document]:
    transcript = YouTubeTranscriptApi.get_transcript(id)
    with YoutubeDL() as ydl:
        info_dict = ydl.extract_info(id, download=False)
        title = info_dict.get('title', None)
    text = ""
    for i in transcript:
        text += i['text'] + " "
    documents = [Document(page_content=text, meta_data={"type":"youtube video", "title":title, "video id":id})]
    return documents, title
